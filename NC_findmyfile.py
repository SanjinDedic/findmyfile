import os
import sys
import fitz
from docx import Document
from pptx import Presentation
import argparse
import pandas as pd
class File:
    
    def __init__(self, path, extension=".txt"):
        self.path = path
        self.name = os.path.basename(self.path)
        self.extension = extension
        self.readable = True
        self.data = ""
    def print(self):
        
        print(self.data)
    def search(self, keyword, keyword2=None):
        
        info = self.data.lower()
        if keyword2:
            return keyword.lower() in info and keyword2.lower() in info
        else:
            return keyword.lower() in info
class PptxFile(File):
    
    def __init__(self, path, extension=".pptx"):
        super().__init__(path, extension)
        self.read()
    def read(self):
        
        try:
            data = ''
            prs = Presentation(self.path)
            values = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            data = ' '.join(values)
            self.data = data
        except:
            self.readable = False
class DocxFile(File):
    
    def __init__(self, path, extension=".docx"):
        super().__init__(path, extension)
        self.read()
    def read(self):
        
        try:
            data = ''
            doc = Document(self.path)
            for para in doc.paragraphs:
                data += '\n' + para.text
            self.data = data
        except:
            self.readable = False
class XlsxFile(File):
    
    def __init__(self, path, extension=".xlsx"):
        super().__init__(path, extension)
        self.read()
    def read(self):
        
        try:
            data = ''
            df = pd.read_excel(self.path)
            data = df.to_string()
            self.data = data
        except:
            self.readable = False
class TxtFile(File):
    
    def __init__(self, path, extension=".txt"):
        super().__init__(path, extension)
        self.read()
    def read(self):
        
        try:
            with open(self.path, 'r', encoding="utf8") as f:
                self.data = f.read()
        except:
            self.readable = False
class PdfFile(File):
    
    def __init__(self, path, extension=".pdf"):
        super().__init__(path, extension)
        self.read()
    def read(self):
        
        try:
            doc = fitz.open(self.path)
            text = ""
            for page in doc:
                text += page.get_text()
            self.data = text
        except:
            self.readable = False
class ProgressBar():
    def __init__(self, total, prefix='', suffix='', decimals=1, length=50, fill='█', print_end="\r"):
        self.total = total
        self.prefix = prefix
        self.suffix = suffix
        self.decimals = decimals
        self.length = length
        self.fill = fill
        self.print_end = print_end
        self.iteration = 0
    def print_progress(self, iteration):
        self.iteration = iteration
        percent = ("{0:." + str(self.decimals) + "f}").format(100 *
                                                              (self.iteration / float(self.total)))
        filled_length = int(self.length * self.iteration // self.total)
        bar = self.fill * filled_length + '-' * (self.length - filled_length)
        print(f'\r{self.prefix} |{bar}| {percent}% {self.suffix}',
              end=self.print_end)
        if self.iteration == self.total:
            print()
class FilesDB():
    def __init__(self, path=os.getcwd()):
        self.path = path
        self.files = []
        self.add_file_data()
    def count_files(self):
        
        file_count = 0
        for root, dirs, files in os.walk(self.path, topdown=False):
            for file in files:
                if file.endswith((".pptx", ".docx", ".xlsx", ".pdf", ".txt")):
                    file_count += 1
        return file_count
    def add_file_data(self):
        
        total_files = self.count_files()
        progress = ProgressBar(
            total_files, prefix='Loading Files:', suffix='Complete', length=50)
        loaded_files = 0
        readers = {'pptx': PptxFile,
                   'docx': DocxFile,
                   'xlsx': XlsxFile,
                   'pdf': PdfFile,
                   'txt': TxtFile
                   }
        for root, dirs, files in os.walk(self.path, topdown=False):
            for file in files:
                ext = os.path.splitext(file)[-1].lstrip('.')
                if ext in readers:
                    data_file = readers[ext](os.path.join(root, file))
                    self.files.append(data_file)
                    loaded_files += 1
                    progress.print_progress(loaded_files)
    def search(self, keyword, keyword2=None):
        counter = 0
        if keyword2:
            print("Files matching keywords: ", keyword, keyword2)
        else:
            print("Files matching keyword: ", keyword)
        for file in self.files:
            counter += 1
            if file.search(keyword, keyword2):
                print(file.name)
        print("\nFiles not readable:")
        for file in self.files:
            if file.readable == False:
                print(file.name)
        print('Files searched:', counter, '\n')
if __name__ == "__main__":
    db = FilesDB()
    parser = argparse.ArgumentParser(prog="findmyfile")
    parser.add_argument('-path', '--source_path', type=str, default=os.getcwd(),
                        help="Takes the directory path in quotes as input")
    parser.add_argument('-keyword1', '--text_pattern',
                        type=str, help="Takes the text pattern as input")
    parser.add_argument('-keyword2', '--text_pattern2', type=str,
                        nargs='?', help="Takes the second text pattern as input")
    parser.add_argument('save', nargs='?', help="Saves a log file in a folder")
    parser.add_argument('verbose', nargs='?',
                        help="Shows the text found in the file")
    args = parser.parse_args()
    keyword1 = args.text_pattern
    keyword2 = args.text_pattern2
    db.search(keyword1, keyword2)
