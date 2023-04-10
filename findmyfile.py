import os
import sys
import fitz
from docx import Document
from pptx import Presentation
import argparse
import pandas as pd


class File:
    """
    The base File class represents a generic file object. It holds the path, name, extension and data of the file.
    It also provides methods to print the file data and search for keywords within the file data.
    """
    def __init__(self, path, extension=".txt"):
        self.path = path
        self.name = os.path.basename(self.path)
        self.extension = extension
        self.readable = True
        self.data = ""

    def print(self):
        """Prints the file data."""
        print(self.data)

    def search(self, keyword, keyword2=None):
        """
        Searches the file data for the given keyword(s). Returns True if the keyword(s) are found,
        otherwise returns False.
        """
        info = self.data.lower()
        if keyword2:
            return keyword.lower() in info and keyword2.lower() in info
        else:
            return keyword.lower() in info


class PptxFile(File):
    """
    The PptxFile class represents a PowerPoint (.pptx) file. It inherits from the File class and
    overrides the read() method to read PowerPoint files.
    """
    def __init__(self, path, extension=".pptx"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        """Reads the PowerPoint file and extracts the text."""
        try:
            data=''
            prs = Presentation(self.path)
            values=[shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            data=' '.join(values)
            self.data = data
        except:
            self.readable = False


class DocxFile(File):
    """
    The DocxFile class represents a Word (.docx) file. It inherits from the File class and
    overrides the read() method to read Word files.
    """
    def __init__(self, path, extension=".docx"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        """Reads the Word file and extracts the text."""
        try:
            data=''
            doc = Document(self.path)
            for para in doc.paragraphs:
                data += '\n'+ para.text
            self.data = data
        except:
            self.readable = False


class XlsxFile(File):
    """
    The XlsxFile class represents an Excel (.xlsx) file. It inherits from the File class and
    overrides the read() method to read Excel files.
    """
    def __init__(self, path, extension=".xlsx"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        """Reads the Excel file and converts the data to a string."""
        try:
            data=''
            df = pd.read_excel(self.path)
            data = df.to_string()
            self.data = data
        except:
            self.readable = False


class TxtFile(File):
    """
    The TxtFile class represents a text (.txt) file. It inherits from the File class and
    overrides the read() method to read text files.
    """
    def __init__(self, path, extension=".txt"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        """Reads the text file and stores the data as a string."""
        try:
            with open(self.path, 'r', encoding="utf8") as f:
                self.data = f.read()
        except:
            self.readable = False


class PdfFile(File):
    """
    The PdfFile class represents a PDF (.pdf) file. It inherits from the File class and
    overrides the read() method to read PDF files.
    """
    def __init__(self, path, extension=".pdf"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        """Reads the PDF file and extracts the text."""
        try:
            doc = fitz.open(self.path)
            text = ""
            for page in doc:
                text += page.get_text()
            self.data = text
        except:
            self.readable = False

class ProgressBar():
    def __init__(self, total, prefix='', suffix='', decimals=1, length=50, fill='█', print_end="\r"):
        self.total = total
        self.prefix = prefix
        self.suffix = suffix
        self.decimals = decimals
        self.length = length
        self.fill = fill
        self.print_end = print_end
        self.iteration = 0

    def print_progress(self, iteration):
        self.iteration = iteration
        percent = ("{0:." + str(self.decimals) + "f}").format(100 * (self.iteration / float(self.total)))
        filled_length = int(self.length * self.iteration // self.total)
        bar = self.fill * filled_length + '-' * (self.length - filled_length)
        print(f'\r{self.prefix} |{bar}| {percent}% {self.suffix}', end=self.print_end)
        if self.iteration == self.total:
            print()


class FilesDB():
    def __init__(self, path=os.getcwd()):
        self.path = path
        self.files = []
        self.add_file_data()

    def count_files(self):
        """
        Counts the total number of supported files in the given path.
        """
        file_count = 0
        for root, dirs, files in os.walk(self.path, topdown=False):
            for file in files:
                if file.endswith(".pptx") or file.endswith(".docx") or file.endswith(".xlsx") or file.endswith(".pdf") or file.endswith(".txt"):
                    file_count += 1
        return file_count


    def add_file_data(self):
        """
        Adds file data to the database by creating appropriate file objects for each file type.
        """
        total_files = self.count_files()
        progress = ProgressBar(total_files, prefix='Loading Files:', suffix='Complete', length=50)
        loaded_files = 0

        for root, dirs, files in os.walk(self.path, topdown=False):
            for file in files:
                if file.endswith(".pptx"):
                    pptx_file = PptxFile(os.path.join(root, file))
                    self.files.append(pptx_file)
                    loaded_files += 1
                    progress.print_progress(loaded_files)

                if file.endswith(".docx"):
                    docx_file = DocxFile(os.path.join(root, file))
                    self.files.append(docx_file)
                    loaded_files += 1
                    progress.print_progress(loaded_files)

                if file.endswith(".xlsx"):
                    docx_file = XlsxFile(os.path.join(root, file))
                    self.files.append(docx_file)
                    loaded_files += 1
                    progress.print_progress(loaded_files)

                if file.endswith(".pdf"):
                    txt_file = PdfFile(os.path.join(root, file))
                    self.files.append(txt_file)
                    loaded_files += 1
                    progress.print_progress(loaded_files)

                if file.endswith(".txt"):
                    txt_file = TxtFile(os.path.join(root, file))
                    self.files.append(txt_file)
                    loaded_files += 1
                    progress.print_progress(loaded_files)
                

    def search(self, keyword, keyword2=None):
        counter = 0
        if keyword2:
            print("Files matching keywords: ", keyword, keyword2)
        else:
            print("Files matching keyword: ", keyword)
        for file in self.files:
            counter += 1
            if file.search(keyword, keyword2):
                print(file.name)
        print("\nFiles not readable:")
        for file in self.files:
            if file.readable == False:
                print(file.name)
        print('Files searched:', counter,'\n' )


    def get_matching_files(self, keyword, keyword2=None):
        results = []
        for file in self.files:
            if file.search(keyword, keyword2):
                results.append(file.path + file.name)
        return results

if __name__ == "__main__":
    db = FilesDB()
    parser = argparse.ArgumentParser(prog="findmyfile")
    parser.add_argument('-path', '--source_path',type=str,default=os.getcwd(), help="Takes the directory path in quotes as input")
    parser.add_argument('-keyword1', '--text_pattern',type=str, help="Takes the text pattern as input")
    parser.add_argument('-keyword2', '--text_pattern2',type=str, nargs='?', help="Takes the second text pattern as input")
    parser.add_argument('save', nargs='?', help="Saves a log file in a folder")
    parser.add_argument('verbose', nargs='?', help="Shows the text found in the file")
    args = parser.parse_args()
    keyword1 = args.text_pattern
    keyword2 = args.text_pattern2
    db.search(keyword1, keyword2)
