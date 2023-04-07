import os
import sys
import fitz
from docx import Document
from pptx import Presentation
import argparse
import pandas as pd
import numpy as np
 

class File:
    def __init__(self, path, extension=".txt"):
        self.path = path
        self.name = os.path.basename(self.path)
        self.extension = extension
        self.readable = True
        self.data = ""

    def print(self):
        print(self.data)

    def search(self, keyword, keyword2=None):
        info = self.data.lower()
        if keyword2:
            return keyword.lower() in info and keyword2.lower() in info
        else:
            return keyword.lower() in info


class PptxFile(File):
    def __init__(self, path, extension=".pptx"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        try:
            data=''
            prs = Presentation(self.path)
            values=[shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            data=' '.join(values)
            self.data = data
        except:
            self.readable = False


class DocxFile(File):
    def __init__(self, path, extension=".docx"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        try:
            data=''
            doc = Document(self.path)
            for para in doc.paragraphs:
                data += '\n'+ para.text
            self.data = data
        except:
            self.readable = False


class XlsxFile(File):
    def __init__(self, path, extension=".xlsx"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        try:
            df = pd.read_excel(self.path, engine='openpyxl')
            self.data = df.to_string()
        except:
            self.readable = False


class TxtFile(File):
    def __init__(self, path, extension=".txt"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        try:
            with open(self.path, 'r', encoding="utf8") as f:
                self.data=f.read()
        except:
            self.readable = False


class PdfFile(File):
    def __init__(self, path, extension=".pdf"):
        super().__init__(path, extension)
        self.read()

    def read(self):
        try:
            doc = fitz.open(self.path)
            text = ""
            for page in doc:
                text += page.get_text()
            self.data = text
        except:
            self.readable = False


class FilesDB():
    def __init__(self, path=os.getcwd()):
        self.path = path
        self.files = []
        self.add_file_data()

    def add_file_data(self):
        for root, dirs, files in os.walk(self.path, topdown=False):
            for file in files:
                if file.endswith(".pptx"):
                    pptx_file = PptxFile(os.path.join(root, file))
                    self.files.append(pptx_file) #will it be faster if I instantiate inside the append()?
                if file.endswith(".docx"):
                    docx_file = DocxFile(os.path.join(root, file))
                    self.files.append(docx_file)
                if file.endswith(".xlsx"):
                    docx_file = XlsxFile(os.path.join(root, file))
                    self.files.append(docx_file)
                if file.endswith(".pdf"):
                    txt_file = PdfFile(os.path.join(root, file))
                    self.files.append(txt_file)
                if file.endswith(".txt"):
                    txt_file = TxtFile(os.path.join(root, file))
                    self.files.append(txt_file)
                

    def search(self, keyword, keyword2=None):
        counter = 0
        if keyword2:
            print("Files matching keywords: ", keyword, keyword2)
        else:
            print("Files matching keyword: ", keyword)
        for file in self.files:
            counter += 1
            if file.search(keyword, keyword2):
                print(file.name)
        print("\nFiles not readable:")
        for file in self.files:
            if file.readable == False:
                print(file.name)
        print('Files searched:', counter,'\n' )


    def get_matching_files(self, keyword, keyword2=None):
        results = []
        for file in self.files:
            if file.search(keyword, keyword2):
                results.append(file.path + file.name)
        return results


if __name__ == "__main__":
    db = FilesDB()
    parser = argparse.ArgumentParser(prog="findmyfile")
    parser.add_argument('-path', '--source_path',type=str,default=os.getcwd(), help="Takes the directory path in quotes as input")
    parser.add_argument('-keyword1', '--text_pattern',type=str, help="Takes the text pattern as input")
    parser.add_argument('-keyword2', '--text_pattern2',type=str, nargs='?', help="Takes the second text pattern as input")
    parser.add_argument('save', nargs='?', help="Saves a log file in a folder")
    parser.add_argument('verbose', nargs='?', help="Shows the text found in the file")
    args = parser.parse_args()
    keyword1 = args.text_pattern
    keyword2 = args.text_pattern2
    db.search(keyword1, keyword2)