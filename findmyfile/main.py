import os
import fitz
from docx import Document
from pptx import Presentation
import argparse
import pandas as pd


class File:
    """
    The base File class represents a generic file object. It holds the path, name, extension and data of the file.
    It also provides methods to print the file data and search for keywords within the file data.
    """

    def __init__(self, path, extension=".txt"):
        self.path = path
        self.name = os.path.basename(self.path)
        self.extension = extension
        self.readable = True
        self.data = ""

    def print(self, chars=100):
        """Prints the file data."""
        print("------------------------------------------------")
        print(self.data[:chars])
        print("------------------------------------------------")

    def search_all(self, *keywords):
        """
        Searches the file data for the given keyword(s). Returns True if all the keyword(s) are found,
        otherwise returns False.
        """
        info = self.data.lower()
        return all(keyword.lower() in info for keyword in keywords)

    def search_any(self, *keywords):
        """
        Searches the file data for the given keyword(s). Returns True if any of the keyword(s) are found,
        otherwise returns False.
        """
        info = self.data.lower()
        return any(keyword.lower() in info for keyword in keywords)

    def read(self):
        pass  # This would be overridden by each subclass


class PptxFile(File):
    """
    The PptxFile class represents a PowerPoint (.pptx) file. It inherits from the File class and
    overrides the read() method to read PowerPoint files.
    """

    def read(self):
        """Reads the PowerPoint file and extracts the text."""
        try:
            data = ''
            prs = Presentation(self.path)
            values = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(
                shape, "text")]
            data = ' '.join(values)
            self.data = data
        except Exception as e:
            print(f"Failed to read PptxFile: {self.path}. Error: {e}")
            self.readable = False


class DocxFile(File):
    """
    The DocxFile class represents a Word (.docx) file. It inherits from the File class and
    overrides the read() method to read Word files.
    """

    def read(self):
        """Reads the Word file and extracts the text."""
        try:
            data = ''
            doc = Document(self.path)
            for para in doc.paragraphs:
                data += '\n' + para.text
            self.data = data
        except Exception as e:
            print(f"Failed to read Docx file: {self.path}. Error: {e}")
            self.readable = False


class XlsxFile(File):
    """
    The XlsxFile class represents an Excel (.xlsx) file. It inherits from the File class and
    overrides the read() method to read Excel files.
    """

    def read(self):
        """Reads the Excel file and converts the data to a string."""
        try:
            data = ''
            df = pd.read_excel(self.path)
            data = df.to_string()
            self.data = data
        except Exception as e:
            print(f"Failed to read Excel file: {self.path}. Error: {e}")
            self.readable = False

class TxtFile(File):
    """
    The TxtFile class represents a text (.txt) file. It inherits from the File class and
    overrides the read() method to read text files.
    """

    def read(self):
        """Reads the text file and stores the data as a string."""
        try:
            with open(self.path, 'r', encoding="utf8") as f:
                self.data = f.read()
        except:
            self.readable = False


class PyFile(File):
    """
    The PyFile class represents a Python (.py) file. It inherits from the File class and
    overrides the read() method to read Python files.
    """

    def read(self):
        """Reads the Python file and stores the data as a string."""
        try:
            with open(self.path, 'r', encoding="utf8") as f:
                self.data = f.read()
        except Exception as e:
            print(f"Failed to read PyFile: {self.path}. Error: {e}")
            self.readable = False


class PdfFile(File):
    """
    The PdfFile class represents a PDF (.pdf) file. It inherits from the File class and
    overrides the read() method to read PDF files.
    """

    def read(self):
        """Reads the PDF file and extracts the text."""
        try:
            doc = fitz.open(self.path)
            text = ""
            for page in doc:
                text += page.get_text()
            self.data = text
        except Exception as e:
            print(f"Failed to read PdfFile: {self.path}. Error: {e}")
            self.readable = False


class ProgressBar():
    """
    This class represents a progress bar. It is used to display the progress of a search task.
    A progress bar is created every time a search task is started and it represents the loading
    of all files into a FilesDB object.
    """

    def __init__(self, total, prefix='', suffix='', decimals=1, length=50, fill='█', print_end="\r"):
        self.total = total
        self.prefix = prefix
        self.suffix = suffix
        self.decimals = decimals
        self.length = length
        self.fill = fill
        self.print_end = print_end
        self.iteration = 0

    def print_progress(self, iteration):
        # detailed docstring
        """Prints the progress bar to the console using the given iteration."""
        self.iteration = iteration
        percent = ("{0:." + str(self.decimals) + "f}").format(100 *
                                                              (self.iteration / float(self.total)))
        filled_length = int(self.length * self.iteration // self.total)
        bar = self.fill * filled_length + '-' * (self.length - filled_length)
        print(f'\r{self.prefix} |{bar}| {percent}% {self.suffix}',
              end=self.print_end)
        if self.iteration == self.total:
            print()


class FilesDB():
    """
    The FilesDB class represents a database of files. It has the add_file_data() method which adds all
    files in a given path to the database. Each file has object has a name, content, extension and path.
    The search method allows users to search for files in the database by keyword(s).
    """

    def __init__(self, path=os.getcwd()):
        self.path = path
        self.files = []
        self.add_file_data()

    def count_files(self):
        """
        Counts the total number of supported files in the given path.
        """
        file_count = 0
        for root, dirs, files in os.walk(self.path, topdown=False):
            for file in files:
                if file.endswith((".pptx", ".docx", ".xlsx", ".pdf", ".txt", ".py")):
                    file_count += 1
        return file_count

    def add_file_data(self):
        """
        Adds file data to the database by creating appropriate file objects for each file type.
        """
        total_files = self.count_files()
        progress = ProgressBar(
            total_files, prefix='Loading Files:', suffix='Complete', length=50)
        loaded_files = 0

        readers = {'pptx': PptxFile,
                   'docx': DocxFile,
                   'xlsx': XlsxFile,
                   'pdf': PdfFile,
                   'txt': TxtFile,
                   'py': PyFile
                   }

        for root, dirs, files in os.walk(self.path, topdown=False):
            for file in files:
                ext = os.path.splitext(file)[-1].lstrip('.')
                if ext in readers:
                    # creates a file object
                    db_file = readers[ext](
                        os.path.join(root, file), extension=ext)
                    db_file.read()  # reads file data into the object
                    self.files.append(db_file)
                    loaded_files += 1
                    progress.print_progress(loaded_files)

    def search(self, print_ans=False, *keywords):
        """Searches for files in the database by keyword(s). Using the print_ans argument, the user can choose
        to print the results to the console or return a dictionary of the results."""
        if print_ans:
            counter = 0
            print("Files matching keywords: ", ', '.join(keywords))
            for file in self.files:
                counter += 1
                if file.search_all(*keywords):
                    #print name up to 30 chars then at 32 char print '|path: ' followed by path
                    print(f'{file.name[:36]:<36} | path: {file.path}')
            print("\nFiles not readable:")
            for file in self.files:
                if file.readable == False:
                    print(f'{file.name[:36]:<36} | path: {file.path}')
            print('Files searched:', counter, '\n')
        else:
            # returns a dictionary where key is file.name and value is file.data
            return {file.name: file.data for file in self.files if file.search_all(*keywords)}

def main():
    # Add your main script here
    db = FilesDB()
    parser = argparse.ArgumentParser(prog="fmf")
    parser.add_argument('-path', '--source_path', type=str, default=os.getcwd(),
                        help="Takes the directory path in quotes as input")
    parser.add_argument('-w', '--words', nargs='*', default=[],
                        help="Takes any number of text patterns as input. Multiple keywords should be separated by space.")
    parser.add_argument('save', nargs='?', help="Saves a log file in a folder")
    parser.add_argument('verbose', nargs='?',
                        help="Shows the text found in the file")
    args = parser.parse_args()
    db.search(*args.words)
